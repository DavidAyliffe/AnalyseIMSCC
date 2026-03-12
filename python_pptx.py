#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
python_pptx.py - Extract plain text from a PowerPoint (.pptx) file

Usage:
    python python_pptx.py <input.pptx> <output.txt>

Arguments:
    input.pptx   Path to the source PowerPoint file
    output.txt   Path for the output plain-text file (created or overwritten)

Output format:
    For each slide, writes a separator block followed by the slide ID and then
    the text content of every shape on that slide.  Consecutive blank lines are
    collapsed to a single blank line and consecutive spaces are collapsed too.

Called by AnalyseIMSCC.pl via:
    system("/usr/local/bin/python3.11", "python_pptx.py", pptxfile, txtfile)

Dependencies:
    python-pptx  (pip install python-pptx)
"""

import sys
import re
from pptx import Presentation


def extract_text(pptx_path, output_path):
    """
    Extract all text from a .pptx file and write it to a plain-text file.

    Each slide is preceded by a separator line and the slide ID.
    Consecutive blank lines and spaces are normalised to single occurrences.

    Parameters
    ----------
    pptx_path   : str  Path to the input .pptx file
    output_path : str  Path to the output .txt file (overwritten if it exists)
    """
    prs = Presentation(pptx_path)

    # Open in write mode ('w') so we always start with a fresh file.
    # Previously this used append mode ('a') which caused old content to
    # accumulate across repeated runs.
    with open(output_path, "w", encoding="utf-8") as f:
        for slide in prs.slides:
            # Write a separator and the slide identifier
            f.write("----------------------\n")
            f.write(str(slide.slide_id) + "\n")
            f.write("----------------------\n")

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                slide_text = shape.text.strip()

                # Collapse multiple consecutive blank lines into one
                slide_text = re.sub(r'\n{2,}', '\n', slide_text)

                # Collapse multiple consecutive spaces into one
                slide_text = re.sub(r' {2,}', ' ', slide_text)

                # Ensure each shape's text is separated from the next
                slide_text += "\n\n"

                f.write(slide_text)


def main():
    """
    Entry point: validate arguments then call extract_text().
    """
    if len(sys.argv) != 3:
        print(
            "Usage: python_pptx.py <input.pptx> <output.txt>",
            file=sys.stderr,
        )
        sys.exit(2)

    pptx_path   = sys.argv[1]
    output_path = sys.argv[2]

    extract_text(pptx_path, output_path)


if __name__ == "__main__":
    main()
